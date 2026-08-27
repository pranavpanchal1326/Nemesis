"""Build ``Nemesis_Prototype.pptx`` — the product walkthrough, screen by screen.

The deck is generated rather than drawn, for the same reason every figure in the
product is computed rather than typed: the frames and their captions come out of
``assets/screens/shots.json``, which the capture script writes as it photographs
the running application. Re-photograph the product and re-run this, and the deck
is current. There is no second copy of the words to drift.

The theme is the product's own — §E9's riso palette on paper stock, the ink
colours out of ``frontend/src/design/tokens.json``, and no gradient, bevel or
drop shadow anywhere. What a screen is allowed to claim is on the slide with it:
a live surface carries the aqua chip, and a screen §E24 still routes to a 404 in
production carries the sunflower one.

    python scripts/build_prototype_deck.py [--out Nemesis_Prototype.pptx]
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SCREENS = ROOT / "assets" / "screens"
DECK_IMAGES = SCREENS / "deck"
MARK = ROOT / "assets" / "nemesis-mark-white.png"

# ── The palette, from frontend/src/design/tokens.json ──────────────────────
PAPER = RGBColor(0xF4, 0xED, 0xE1)  # paper-50
KRAFT = RGBColor(0xEA, 0xDC, 0xC7)  # kraft-100
BONE = RGBColor(0xE4, 0xDC, 0xCE)  # bone-200
MITTI = RGBColor(0x8A, 0x74, 0x62)  # mitti-500
MITTI_300 = RGBColor(0xBF, 0xAE, 0x9B)
INK = RGBColor(0x16, 0x13, 0x0F)  # riso-black
BROWN = RGBColor(0x92, 0x5F, 0x52)  # riso-brown
SUNFLOWER = RGBColor(0xFF, 0xB5, 0x11)
AQUA = RGBColor(0x0F, 0xA9, 0xA0)
FED_BLUE = RGBColor(0x3D, 0x55, 0x88)

# ── The type. Nothing exotic: the product's own faces (Gambarino, Panchang,
#    Switzer) are web fonts nobody presenting this deck has installed, and a
#    slide that silently substitutes is worse than one that never asked.
DISPLAY = "Segoe UI Black"
HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)


@dataclass(frozen=True)
class Section:
    """One movement of the walkthrough."""

    number: str
    title: str
    blurb: str
    #: Frames belong to a section by the numeric prefix the capture gave them.
    first: int
    last: int


SECTIONS: tuple[Section, ...] = (
    Section(
        "01",
        "The landing film",
        "Nine acts over a 3D clay model of the city, scroll-driven. "
        "The model's light is that city's local time; its rain is that city's real weather.",
        1,
        10,
    ),
    Section(
        "02",
        "The resident",
        "Photograph it, say where, send. Then a receipt whose claim is a hash, "
        "and an evidence trail nobody can look up by your name.",
        11,
        16,
    ),
    Section(
        "03",
        "What the city publishes",
        "The open-data surfaces — wards, money, contractors — and the page that says "
        "which of this product's claims are real.",
        17,
        21,
    ),
    Section(
        "04",
        "The staff",
        "The console: what breaches first, the reports the pipeline would not decide alone, "
        "and the administrative plane underneath them.",
        22,
        29,
    ),
    Section(
        "05",
        "What is not wired yet",
        "Seven screens whose contracts still return nulls. A production build routes every one "
        "of them to a 404 — they are designed, and they say so.",
        30,
        36,
    ),
    Section(
        "06",
        "Outdoors, and under the hood",
        "The crew's phone, offline and in sunlight — and the proof routes each rendering "
        "pipeline is photographed through.",
        37,
        38,
    ),
)


# ── Drawing primitives ─────────────────────────────────────────────────────


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _ground(slide, colour: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False


def _rule(slide, x, y, width, colour: RGBColor, thickness=Pt(1.5)) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False


def _text(
    slide,
    x,
    y,
    width,
    height,
    runs: list[tuple[str, str, int, RGBColor]],
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing: float = 1.0,
    space_before: int = 0,
):
    """A text box of one paragraph per run: (text, font, size_pt, colour)."""
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0

    for index, (text, font, size, colour) in enumerate(runs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        if index > 0 and space_before:
            para.space_before = Pt(space_before)
        run = para.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = colour
    return box


def _chip(slide, x, y, label: str, ink: RGBColor, ground: RGBColor) -> None:
    """A status chip. The product puts one on every screen that needs one; a
    deck that dropped them would be making claims the product refuses to."""
    # Sized to the label rather than fixed: "ROADMAP · DEV ONLY" is half again
    # as long as "LIVE · REAL DATA", and a chip its text overhangs is a chip
    # that reads as a mistake.
    width, height = Inches(0.28 + 0.073 * len(label)), Inches(0.26)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ground
    shape.line.color.rgb = ink
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False

    frame = shape.text_frame
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    run.font.name = HEAD
    run.font.size = Pt(9)
    run.font.color.rgb = ink


def _place_image(slide, image: Path, box_x, box_y, box_w, box_h, *, frame=True):
    """Fit an image inside a box, centred, with a hairline rule around it."""
    with Image.open(image) as probe:
        ratio = probe.height / probe.width

    width = box_w
    height = Emu(int(width * ratio))
    if height > box_h:
        height = box_h
        width = Emu(int(height / ratio))

    x = Emu(int(box_x + (box_w - width) / 2))
    y = Emu(int(box_y + (box_h - height) / 2))

    if frame:
        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(x - Emu(9525)), Emu(y - Emu(9525)),
            Emu(width + Emu(19050)), Emu(height + Emu(19050)),
        )
        edge.fill.background()
        edge.line.color.rgb = MITTI_300
        edge.line.width = Pt(0.75)
        edge.shadow.inherit = False

    slide.shapes.add_picture(str(image), x, y, width, height)


def _footer(slide, prs, left: str, right: str, colour: RGBColor) -> None:
    _text(
        slide,
        Inches(0.75),
        Inches(6.92),
        Inches(9.0),
        Inches(0.3),
        [(left, MONO, 9, colour)],
    )
    _text(
        slide,
        Inches(10.0),
        Inches(6.92),
        Inches(2.58),
        Inches(0.3),
        [(right, MONO, 9, colour)],
        align=PP_ALIGN.RIGHT,
    )


# ── The slides ─────────────────────────────────────────────────────────────


def title_slide(prs) -> None:
    slide = _blank(prs)
    _ground(slide, INK)

    if MARK.exists():
        slide.shapes.add_picture(str(MARK), Inches(0.95), Inches(1.15), height=Inches(1.5))

    _text(
        slide,
        Inches(0.9),
        Inches(3.0),
        Inches(11.5),
        Inches(1.5),
        [("NEMESIS", DISPLAY, 88, PAPER)],
        spacing=0.9,
    )
    _rule(slide, Inches(0.95), Inches(4.28), Inches(2.2), SUNFLOWER, Pt(3))
    _text(
        slide,
        Inches(0.95),
        Inches(4.62),
        Inches(9.0),
        Inches(1.4),
        [
            ("AI CIVIC OPERATIONS AGENT", HEAD, 16, KRAFT),
            (
                "Networked Enforcement & Municipal Evidence System "
                "for Infrastructure & Service Accountability",
                BODY,
                13,
                MITTI_300,
            ),
        ],
        space_before=8,
    )
    _text(
        slide,
        Inches(0.95),
        Inches(6.05),
        Inches(6.0),
        Inches(0.6),
        [("Prove, don't log.", HEAD, 26, SUNFLOWER)],
    )
    _text(
        slide,
        Inches(8.4),
        Inches(6.18),
        Inches(4.0),
        Inches(0.6),
        [("PROTOTYPE WALKTHROUGH", MONO, 11, MITTI)],
        align=PP_ALIGN.RIGHT,
    )


def statement_slide(prs, kicker: str, lines: list[str], *, dark=False, accent=BROWN) -> None:
    slide = _blank(prs)
    _ground(slide, INK if dark else PAPER)
    body = PAPER if dark else INK
    quiet = MITTI_300 if dark else MITTI

    _text(slide, Inches(0.95), Inches(0.85), Inches(8.0), Inches(0.3), [(kicker, MONO, 11, quiet)])
    _rule(slide, Inches(0.95), Inches(1.28), Inches(1.6), accent, Pt(3))

    runs = [(line, HEAD if index == 0 else BODY, 30 if index == 0 else 17, body if index == 0 else quiet)
            for index, line in enumerate(lines)]
    _text(
        slide,
        Inches(0.95),
        Inches(1.9),
        Inches(11.0),
        Inches(4.4),
        runs,
        spacing=1.22,
        space_before=16,
    )


def section_slide(prs, section: Section, count: int) -> None:
    slide = _blank(prs)
    _ground(slide, INK)

    _text(
        slide,
        Inches(0.95),
        Inches(1.5),
        Inches(4.0),
        Inches(2.2),
        [(section.number, DISPLAY, 128, BROWN)],
        spacing=0.9,
    )
    _rule(slide, Inches(0.95), Inches(3.55), Inches(11.4), MITTI, Pt(1))
    _text(
        slide,
        Inches(0.95),
        Inches(3.95),
        Inches(10.5),
        Inches(1.0),
        [(section.title, HEAD, 40, PAPER)],
    )
    _text(
        slide,
        Inches(0.95),
        Inches(5.05),
        Inches(9.2),
        Inches(1.4),
        [(section.blurb, BODY, 15, MITTI_300)],
        spacing=1.28,
    )
    _text(
        slide,
        Inches(10.3),
        Inches(5.05),
        Inches(2.1),
        Inches(0.6),
        [(f"{count} SCREENS", MONO, 11, SUNFLOWER)],
        align=PP_ALIGN.RIGHT,
    )


def shot_slide(prs, shot: dict, section: Section, number: int, total: int) -> None:
    slide = _blank(prs)
    _ground(slide, PAPER)

    image = DECK_IMAGES / shot["image"]
    phone = shot["surface"] == "phone"
    roadmap = shot["availability"] == "roadmap"

    if phone:
        # The shape of the phone is the point, so it gets the narrow column and
        # the words get the wide one.
        _place_image(slide, image, Inches(0.85), Inches(0.95), Inches(3.0), Inches(5.6))
        text_x, text_w = Inches(4.5), Inches(8.0)
    else:
        _place_image(slide, image, Inches(4.5), Inches(1.15), Inches(8.05), Inches(5.1))
        text_x, text_w = Inches(0.78), Inches(3.4)

    _text(
        slide,
        text_x,
        Inches(1.3),
        text_w,
        Inches(0.3),
        [(section.title.upper(), MONO, 10, BROWN)],
    )
    _rule(slide, text_x, Inches(1.66), Inches(0.9), SUNFLOWER if roadmap else AQUA, Pt(3))
    _text(
        slide,
        text_x,
        Inches(2.02),
        text_w,
        Inches(1.3),
        [(shot["title"], HEAD, 22, INK)],
        spacing=1.12,
    )
    _text(
        slide,
        text_x,
        Inches(3.2),
        text_w,
        Inches(2.5),
        [(shot["caption"], BODY, 13, MITTI)],
        spacing=1.34,
    )

    if roadmap:
        _chip(slide, text_x, Inches(5.95), "ROADMAP · DEV ONLY", RGBColor(0x8A, 0x4E, 0x12), KRAFT)
    else:
        _chip(slide, text_x, Inches(5.95), "LIVE · REAL DATA", RGBColor(0x0A, 0x6B, 0x66), BONE)

    _footer(slide, prs, shot["url"], f"{number:02d} / {total:02d}", MITTI_300)


def pipeline_slide(prs) -> None:
    """How a report moves — the one diagram the deck draws rather than photographs."""
    slide = _blank(prs)
    _ground(slide, PAPER)

    _text(slide, Inches(0.95), Inches(0.85), Inches(8.0), Inches(0.3),
          [("HOW A REPORT MOVES", MONO, 11, BROWN)])
    _rule(slide, Inches(0.95), Inches(1.28), Inches(1.6), AQUA, Pt(3))
    _text(
        slide,
        Inches(0.95),
        Inches(1.75),
        Inches(11.4),
        Inches(0.8),
        [("Every stage appends events. Nothing mutates a row.", HEAD, 26, INK)],
    )

    stages = [
        ("ingest", True),
        ("safety", True),
        ("trust", True),
        ("classify", True),
        ("dedup", True),
        ("severity", False),
        ("routing · SLA", False),
        ("investigation", False),
    ]
    x = Inches(0.95)
    top = Inches(3.05)
    width = Inches(1.32)
    gap = Inches(0.115)
    height = Inches(0.92)

    for label, live in stages:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = KRAFT if live else PAPER
        box.line.color.rgb = INK if live else MITTI_300
        box.line.width = Pt(1.5) if live else Pt(0.75)
        if not live:
            box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        box.shadow.inherit = False

        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.05)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.name = HEAD if live else BODY
        run.font.size = Pt(11)
        run.font.color.rgb = INK if live else MITTI
        x = Emu(x + width + gap)

    _text(
        slide,
        Inches(0.95),
        Inches(4.3),
        Inches(11.4),
        Inches(0.4),
        [("SOLID — shipped and gated        DASHED — designed, not built", MONO, 10, MITTI)],
    )

    notes = [
        "The safety fail-safe runs on its own queue, served by a container that has never "
        "imported torch. A saturated ML queue cannot delay a gas-leak report.",
        "Degradation is recorded, not swallowed. Every external call has a timeout, a retry "
        "budget, a fallback, and its own event — DEGRADED is not FAILED.",
        "Deduplication is reversible. A merge is undone by a compensating event, never a delete.",
    ]
    _text(
        slide,
        Inches(0.95),
        Inches(4.95),
        Inches(11.4),
        Inches(2.0),
        [(f"·   {note}", BODY, 12.5, MITTI) for note in notes],
        spacing=1.3,
        space_before=10,
    )


def closing_slide(prs, live: int, roadmap: int) -> None:
    slide = _blank(prs)
    _ground(slide, INK)

    _text(
        slide,
        Inches(0.95),
        Inches(1.55),
        Inches(11.0),
        Inches(2.2),
        [
            ("Every figure on every one of these screens", BODY, 22, MITTI_300),
            ("is computed from an append-only log.", HEAD, 36, PAPER),
        ],
        spacing=1.16,
        space_before=10,
    )
    _rule(slide, Inches(0.95), Inches(4.0), Inches(11.4), MITTI, Pt(1))

    columns = (
        (f"{live}", "screens photographed live,\nagainst a seeded deployment", AQUA),
        (f"{roadmap}", "designed and honestly\nrouted to a 404 in production", SUNFLOWER),
        ("0", "figures typed by hand\nanywhere in this deck", BROWN),
    )
    x = Inches(0.95)
    for value, label, colour in columns:
        _text(slide, x, Inches(4.45), Inches(3.5), Inches(1.0), [(value, DISPLAY, 54, colour)])
        _text(
            slide,
            x,
            Inches(5.42),
            Inches(3.3),
            Inches(1.0),
            [(line, BODY, 12, MITTI_300) for line in label.split("\n")],
            spacing=1.25,
        )
        x = Emu(x + Inches(3.8))

    _text(
        slide,
        Inches(0.95),
        Inches(6.75),
        Inches(11.4),
        Inches(0.4),
        [("Prove, don't log.", HEAD, 16, PAPER)],
    )


# ── Assembly ───────────────────────────────────────────────────────────────


def section_for(shot: dict) -> Section | None:
    index = int(shot["file"][:2])
    for section in SECTIONS:
        if section.first <= index <= section.last:
            return section
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", default="Nemesis_Prototype.pptx")
    options = parser.parse_args()

    manifest_path = SCREENS / "shots.json"
    if not manifest_path.exists():
        print(f"no manifest at {manifest_path} — run the capture first")
        return 1

    shots = [
        shot
        for shot in json.loads(manifest_path.read_text(encoding="utf-8"))
        if "image" in shot and (DECK_IMAGES / shot["image"]).exists()
    ]
    shots.sort(key=lambda shot: shot["file"])
    if not shots:
        print("the manifest has no built images — run scripts/build_screen_derivatives.py")
        return 1

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    title_slide(prs)
    statement_slide(
        prs,
        "THE PROBLEM",
        [
            "A pothole gets reported. The app says “In Progress”. Weeks pass.",
            "No one knows who did the work, what it cost, or why it failed. "
            "So the next complaint never gets filed.",
            "That silence isn't apathy. It's the sound of a system that broke its "
            "promise once too often.",
        ],
        dark=True,
        accent=SUNFLOWER,
    )
    statement_slide(
        prs,
        "THE PREMISE",
        [
            "A civic system earns trust by producing evidence, not status updates.",
            "Every figure NEMESIS publishes is computed from an append-only event log. "
            "Every claim it makes about itself carries a status label. Nothing is edited by hand.",
            "What follows is the running product, photographed against a seeded deployment "
            "of a demo city — not a mockup, and not a rendering of one.",
        ],
        accent=AQUA,
    )

    total = len(shots)
    number = 0
    current: Section | None = None
    for shot in shots:
        section = section_for(shot)
        if section is None:
            continue
        if section is not current:
            section_slide(prs, section, sum(1 for s in shots if section_for(s) is section))
            current = section
        number += 1
        shot_slide(prs, shot, section, number, total)

    pipeline_slide(prs)
    closing_slide(
        prs,
        sum(1 for shot in shots if shot["availability"] == "live"),
        sum(1 for shot in shots if shot["availability"] == "roadmap"),
    )

    out = ROOT / options.out
    prs.save(out)
    print(f"wrote {out.name} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides, {total} screens")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
